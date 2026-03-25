#!/usr/bin/env python3
"""
generate_pptx.py — Generate a CRE Market Survey PowerPoint from processed listings.

Usage:
    python generate_pptx.py --data ./data/processed_listings.json \
                             --template ./templates/Market_Survey_Template.pptx \
                             --output ./output/Market_Survey.pptx

Slides:
  Slide 1: Title slide
  Slide 2: Map slide (satellite with numbered pins)
  Slide 3+: Detail table rows (Map #, Photo, Address, SF, Year Built, Doors, Price, PPSF, Notes)
"""

import argparse
import json
import os
import re
import sys
import hashlib
import urllib.parse
import urllib.request
from collections import Counter
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from PIL import Image as PILImage

try:
    from generate_map import (
        get_api_key,
        get_mapbox_token,
        fetch_mapbox_satellite_map,
        geocode_listings,
        load_render_config,
        generate_map_image,
        fetch_google_satellite_map,
        generate_google_static_map_from_addresses,
        render_map_fallback,
        get_bounds,
        compute_zoom,
        lat_lon_to_pixel_google,
        compute_pin_layout_positions,
        draw_road_labels,
        draw_white_roads,
        draw_highway_shields,
    )
except ImportError:
    print("WARNING: generate_map.py not found. Map generation will be skipped.")
    def geocode_listings(listings, api_key=None, config=None):
        return [{"map_num": i+1, "address": l.get("address", ""), "lat": None, "lon": None}
                for i, l in enumerate(listings)]
    def get_api_key():
        return None
    def get_mapbox_token():
        return None
    def fetch_mapbox_satellite_map(*a, **kw):
        return False
    def load_render_config(*a, **kw):
        return {}
    def generate_map_image(*a, **kw):
        return {"success": False, "geocoded": [], "valid_count": 0, "used_satellite_base": False, "center_lat": None, "center_lon": None, "zoom": None, "crop_bottom": 0, "render_mode": "unavailable", "qa_warnings": []}
    def fetch_google_satellite_map(*a, **kw):
        return False
    def generate_google_static_map_from_addresses(*a, **kw):
        return False
    def render_map_fallback(geocoded, path, width=1200, height=800):
        pass
    def get_bounds(*a, **kw):
        return {"min_lat": 0, "max_lat": 1, "min_lon": 0, "max_lon": 1}
    def compute_zoom(*a, **kw):
        return 1
    def lat_lon_to_pixel_google(*a, **kw):
        return (0, 0)
    def compute_pin_layout_positions(geocoded, *a, **kw):
        return [{"point": g, "x": 0.0, "y": 0.0} for g in geocoded if g.get("lat") is not None and g.get("lon") is not None]
    def draw_road_labels(*a, **kw):
        pass
    def draw_white_roads(*a, **kw):
        pass
    def draw_highway_shields(*a, **kw):
        pass


# ============================================================
# Known geography / property-type helpers
# ============================================================

KNOWN_STATES = {
    "Montana": "MT", "Wyoming": "WY", "Colorado": "CO", "Idaho": "ID",
    "South Dakota": "SD", "North Dakota": "ND", "Nebraska": "NE", "Utah": "UT",
}

KNOWN_CITIES = {
    "billings", "bozeman", "missoula", "great falls", "helena",
    "casper", "cheyenne", "gillette", "rock springs", "laramie",
    "denver", "colorado springs", "fort collins", "boise", "idaho falls",
}

PROPERTY_TYPE_KEYWORDS = {
    "Industrial": ["industrial", "warehouse", "shop", "bay", "dock", "crane", "mezzanine", "manufacturing", "fabrication"],
    "Commercial": ["office", "retail", "commercial", "storefront", "showroom", "restaurant", "medical"],
    "Land": ["land", "lot", "acreage", "acre", "parcel", "zoned", "vacant"],
    "Residential": ["residential", "house", "home", "apartment", "multi-family"],
}


# ============================================================
# Data loading
# ============================================================

def load_listings(data_path: str) -> list[dict]:
    """Load processed listings JSON."""
    with open(data_path, encoding="utf-8") as f:
        data = json.load(f)
    if isinstance(data, list):
        return data
    if isinstance(data, dict) and "listings" in data:
        return data["listings"]
    raise ValueError(f"Unexpected JSON structure in {data_path}")


def get_field_value(listing: dict, field: str, default: Any = "") -> Any:
    """Extract a field value, handling both flat and nested (value/confidence) formats."""
    val = listing.get(field, default)
    if isinstance(val, dict) and "value" in val:
        return val["value"] or default
    return val if val is not None else default


# ============================================================
# Formatting helpers
# ============================================================

def format_price(price_str: str) -> str:
    if not price_str:
        return "Call for Pricing"
    try:
        num_str = re.sub(r"[^\d.]", "", str(price_str))
        if not num_str:
            return "Call for Pricing"
        num = float(num_str)
        if "/mo" in str(price_str).lower():
            return f"${num:,.0f}/mo"
        return f"${num:,.0f}"
    except (ValueError, TypeError):
        return str(price_str) if price_str else "Call for Pricing"


def format_monthly_price(price_str: str) -> str:
    if not price_str:
        return "Call for Pricing"
    try:
        num_str = re.sub(r"[^\d.]", "", str(price_str))
        if not num_str:
            return "Call for Pricing"
        num = float(num_str)
        return f"${num:,.0f}/mo"
    except (ValueError, TypeError):
        return str(price_str) if price_str else "Call for Pricing"


def format_sf(sf_str: str) -> str:
    if not sf_str:
        return "N/A"
    try:
        num = float(re.sub(r"[^\d.]", "", str(sf_str)))
        return f"{num:,.0f}"
    except (ValueError, TypeError):
        return str(sf_str) if sf_str else "N/A"


def format_sf_cell(sf_str: str) -> str:
    base = format_sf(sf_str)
    return f"{base} SF" if base != "N/A" else "N/A"


def street_only_address(address: str) -> str:
    if not address:
        return ""
    parts = [p.strip() for p in str(address).split(',') if p.strip()]
    if not parts:
        base = str(address).strip()
    else:
        base = parts[0]

    # Clean unit/suite fragments from display text so the slide looks closer
    # to the original survey style. Keep the full address for hyperlink targets.
    base = re.sub(r'\s+#\s*[A-Za-z0-9,\-]+\s*$', '', base).strip()
    base = re.sub(r'\s+(?:Suite|Ste|Unit)\s+[A-Za-z0-9\-]+\s*$', '', base, flags=re.IGNORECASE).strip()
    return base


def google_maps_search_url(address: str) -> str:
    if not address:
        return ""
    from urllib.parse import quote_plus
    return f"https://www.google.com/maps/search/?api=1&query={quote_plus(address)}"


def calculate_ppsf(price_str: str, sf_str: str) -> str:
    try:
        price_clean = re.sub(r"[^\d.]", "", str(price_str or ""))
        sf_clean = re.sub(r"[^\d.]", "", str(sf_str or ""))
        if not price_clean or not sf_clean:
            return "NNN"
        price = float(price_clean)
        sf = float(sf_clean)
        if sf <= 0:
            return "NNN"
        ppsf = price / sf
        return f"${ppsf:.2f}/SF\nNNN"
    except (ValueError, TypeError):
        return "NNN"


def format_notes(notes: str) -> str:
    if not notes:
        return ""
    text = re.sub(r"\s+", " ", str(notes)).strip()
    return text


# ============================================================
# Shape manipulation helpers
# ============================================================

def clear_shape_text(shape) -> None:
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = ""
        if para.runs:
            para.runs[0].text = ""


def set_shape_lines(shape, lines: list, margin_top: int = 0, margin_bottom: int = 0,
                    margin_left: int = 0, margin_right: int = 0,
                    vertical_anchor=None, align=None) -> None:
    """Replace a text shape with explicit lines and consistent in-frame layout."""
    tf = shape.text_frame
    tf.word_wrap = True
    if margin_top:
        tf.margin_top = Emu(margin_top)
    if margin_bottom:
        tf.margin_bottom = Emu(margin_bottom)
    if margin_left:
        tf.margin_left = Emu(margin_left)
    if margin_right:
        tf.margin_right = Emu(margin_right)
    if vertical_anchor is not None:
        tf.vertical_anchor = vertical_anchor
    # Clear existing paragraphs
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)
    for idx, text in enumerate(lines):
        font_size = None
        if isinstance(text, tuple):
            text, font_size = text
        if idx == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        if align is not None:
            para.alignment = align
        run = para.add_run() if not para.runs else para.runs[0]
        run.text = text
        if font_size:
            run.font.size = Pt(font_size)


def remove_shape(shape) -> None:
    sp = shape._element
    sp.getparent().remove(sp)


def get_table_row_elements(table) -> list:
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    return table._tbl.findall(".//a:tr", nsmap)


def delete_table_rows(table, row_indices: list[int]) -> int:
    """Delete table rows by index, excluding the header row at index 0."""
    tr_elements = get_table_row_elements(table)
    deleted = 0
    for row_idx in sorted(row_indices, reverse=True):
        if row_idx == 0 or row_idx >= len(tr_elements):
            continue
        tr_elem = tr_elements[row_idx]
        tr_elem.getparent().remove(tr_elem)
        deleted += 1
    return deleted


def sync_table_shape_height(table_shape) -> None:
    """Resize the table shape height to match the sum of its row heights.

    When rows are deleted from a table, python-pptx leaves the shape's
    height attribute at the original template value.  PowerPoint then
    distributes the leftover space visually, shifting row boundaries
    relative to the XML coordinates we use for floating photo shapes.
    Resetting the shape height locks the rendered rows to their XML positions.
    """
    row_height_sum = sum(row.height for row in table_shape.table.rows)
    table_shape.height = row_height_sum
    print(f"    Table height synced to {row_height_sum / 914400:.4f}\"")


# Font sizes by column index, matching populate_table_row assignments.
_COL_FONT_SIZES = {0: 9, 1: 9, 2: 10.5, 3: 9, 4: 9, 5: 9, 6: 9.5, 7: 8.5, 8: 8.5}


def _estimate_cell_text_height(text: str, col_width_emu: int, font_pt: float) -> int:
    """Return the estimated EMU height a cell needs for its text content.

    Uses a conservative average-character-width for Arial to approximate
    how many wrapped lines PowerPoint will render, then converts to EMU.
    """
    if not text:
        return 0

    # Default cell margins in EMU (python-pptx defaults when unset).
    margin_lr = 2 * 45720  # left + right
    margin_tb = 2 * 45720  # top + bottom

    usable_width_pt = (col_width_emu - margin_lr) / 12700
    # Arial average character width ≈ 50 % of the point size.
    char_width_pt = font_pt * 0.50
    chars_per_line = max(1, int(usable_width_pt / char_width_pt))

    lines = 0
    for paragraph in text.split("\n"):
        paragraph = paragraph.strip()
        if paragraph:
            lines += max(1, -(-len(paragraph) // chars_per_line))  # ceil div
        else:
            lines += 1

    line_height_emu = int(font_pt * 1.25 * 12700)  # ~125 % line spacing
    return lines * line_height_emu + margin_tb


def adjust_row_heights_for_content(table_shape) -> None:
    """Grow data-row heights so the text content fits without PowerPoint auto-expanding.

    PowerPoint treats XML row heights as *minimums* and silently expands rows
    when content overflows.  Floating photo shapes are positioned from the XML
    heights, so any hidden expansion causes a cumulative upward drift of photos.
    Calling this function **after** text is populated and **before** photos are
    placed keeps the two in sync.
    """
    table = table_shape.table
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    trs = table._tbl.findall(".//a:tr", nsmap)

    for row_idx in range(1, len(table.rows)):  # skip header
        row = table.rows[row_idx]
        max_needed = 0

        for col_idx, cell in enumerate(row.cells):
            text = cell.text_frame.text
            if not text.strip():
                continue
            font_pt = _COL_FONT_SIZES.get(col_idx, 9)
            col_width = table.columns[col_idx].width
            needed = _estimate_cell_text_height(text, col_width, font_pt)
            if needed > max_needed:
                max_needed = needed

        if max_needed > row.height:
            old_h = row.height
            row.height = max_needed
            if row_idx < len(trs):
                trs[row_idx].set("h", str(max_needed))
            print(f"    Row {row_idx}: height {old_h / 914400:.3f}\" -> {max_needed / 914400:.3f}\" (text overflow fix)")


# ============================================================
# Metadata detection
# ============================================================

def detect_survey_metadata(listings: list[dict]) -> dict:
    """Detect city, state, and property type from listing data."""
    cities: Counter = Counter()
    states: Counter = Counter()

    for listing in listings:
        address = get_field_value(listing, "address")
        match = re.search(r",\s*([A-Za-z\s]+),\s*([A-Z]{2})\s*\d{5}", address)
        if match:
            city = match.group(1).strip()
            state_abbr = match.group(2).strip()
            cities[city.lower()] += 1
            states[state_abbr] += 1

    city = cities.most_common(1)[0][0].title() if cities else ""
    state_abbr = states.most_common(1)[0][0] if states else ""
    state_full = next((k for k, v in KNOWN_STATES.items() if v == state_abbr), state_abbr)

    # Determine property type
    explicit_types: Counter = Counter()
    for l in listings:
        pt = get_field_value(l, "property_type")
        if pt:
            explicit_types[pt] += 1

    if explicit_types:
        ptype, count = explicit_types.most_common(1)[0]
        confidence = min(0.9, count / len(listings))
    else:
        type_scores: Counter = Counter()
        for l in listings:
            text = " ".join([
                get_field_value(l, "notes", ""),
                get_field_value(l, "source_file", ""),
            ]).lower()
            for ptype, keywords in PROPERTY_TYPE_KEYWORDS.items():
                for kw in keywords:
                    if kw in text:
                        type_scores[ptype] += 1

        total = sum(type_scores.values()) or 1
        ptype = type_scores.most_common(1)[0][0] if type_scores else "Industrial"
        confidence = type_scores.most_common(1)[0][1] / total if type_scores else 0.3

    return {
        "city": city,
        "state": state_full,
        "state_abbr": state_abbr,
        "property_type": ptype,
        "property_type_confidence": confidence,
    }


# ============================================================
# Branding updates (Stage 1)
# ============================================================

def update_title_slide(prs: Presentation, metadata: dict, date_str: str = "") -> None:
    """Update the title slide (slide 1) with dynamic city, property type, and date."""
    slide = prs.slides[0]
    city = metadata.get("city", "")
    state_abbr = metadata.get("state_abbr", "")
    ptype = metadata.get("property_type", "INDUSTRIAL")

    if not date_str:
        date_str = datetime.now().strftime("%B %Y")

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = run.text
                # Replace known city names
                for known_city in KNOWN_CITIES:
                    if known_city.lower() in text.lower():
                        text = re.sub(known_city, city, text, flags=re.IGNORECASE)
                        run.text = text
                # Replace property type
                for old_type in ["INDUSTRIAL", "COMMERCIAL", "LAND", "RESIDENTIAL"]:
                    if old_type in run.text.upper():
                        run.text = run.text.replace(old_type, ptype.upper())
                # Replace market survey date
                match = re.search(
                    r"(January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}",
                    run.text
                )
                if match:
                    run.text = run.text.replace(match.group(0), date_str)

    print(f"  Title slide updated: {city} {ptype.upper()}")


def update_map_banner(prs: Presentation, metadata: dict, map_slide_index: int = 1) -> None:
    """Update the banner text on the map slide."""
    if map_slide_index >= len(prs.slides):
        return
    slide = prs.slides[map_slide_index]
    city = metadata.get("city", "")
    state_abbr = metadata.get("state_abbr", "")
    ptype = metadata.get("property_type", "INDUSTRIAL")

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip().lower()
        is_banner = False
        for known_city in KNOWN_CITIES:
            if known_city in text:
                is_banner = True
                break
        if not is_banner:
            for abbr in KNOWN_STATES.values():
                if abbr.lower() in text:
                    is_banner = True
                    break
        if not is_banner:
            continue

        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run_text = run.text
                for known_city in KNOWN_CITIES:
                    if known_city in run_text.lower():
                        run_text = re.sub(known_city, city, run_text, flags=re.IGNORECASE)
                for abbr in KNOWN_STATES.values():
                    if abbr in run_text.upper():
                        run_text = run_text.replace(abbr, state_abbr)
                for old_type in ["INDUSTRIAL", "COMMERCIAL", "LAND", "RESIDENTIAL"]:
                    if old_type in run_text.upper():
                        run_text = run_text.replace(old_type, ptype.upper())
                run.text = run_text
        print(f"  Map banner updated: {city.upper()}, {state_abbr} ({ptype.upper()})")
        return


def update_detail_headers(prs: Presentation, metadata: dict) -> None:
    """Update 'INDUSTRIAL DETAILS' headers on detail slides to match property type."""
    ptype = metadata.get("property_type", "INDUSTRIAL").upper()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old_type in ["INDUSTRIAL", "COMMERCIAL", "LAND", "RESIDENTIAL"]:
                        if old_type in run.text.upper() and "DETAILS" in run.text.upper():
                            run.text = re.sub(old_type, ptype, run.text, flags=re.IGNORECASE)


# ============================================================
# Map slide cleanup (Stage 3)
# ============================================================

def clean_map_slide(prs: Presentation, map_slide_index: int = 1) -> None:
    """Remove old template pin shapes from the map slide."""
    if map_slide_index >= len(prs.slides):
        return

    slide = prs.slides[map_slide_index]
    shapes_to_remove = []

    for shape in slide.shapes:
        # Keep tables (the property overlay)
        if shape.has_table:
            continue
        # Keep large images (the map background)
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            continue
        # Keep text boxes with substantial text (banner, title)
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if len(text) > 3:
                continue
        # Small shapes are likely old pins
        if shape.width < Inches(1) and shape.height < Inches(1):
            shapes_to_remove.append(shape)
        elif hasattr(shape, 'shape_type') and shape.shape_type == 6:  # GROUP
            if shape.width < Inches(1.5) and shape.height < Inches(1.5):
                shapes_to_remove.append(shape)

    if shapes_to_remove:
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
        print(f"  Removed {len(shapes_to_remove)} old template pin shapes from map slide")
    else:
        print(f"  No old pin shapes found to remove")


# ============================================================
# Table population & pagination (Stage 3)
# ============================================================

def find_table_slides(prs: Presentation) -> dict:
    result = {"map": None, "detail": None}
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_table:
                col_count = len(shape.table.columns)
                if col_count >= 5 and result["detail"] is None:
                    result["detail"] = (i, shape)
                elif col_count < 5 and result["map"] is None:
                    result["map"] = (i, shape)
    return result


def capture_photo_placeholders(slide, table_shape) -> list[dict]:
    """Find legacy picture placeholders so they can be removed before repainting photos."""
    table = table_shape.table
    photo_col_left = table_shape.left + table.columns[0].width
    photo_col_right = photo_col_left + table.columns[1].width
    table_top = table_shape.top
    table_bottom = table_shape.top + table_shape.height

    slots = []
    for shape in slide.shapes:
        if shape.shape_type != 13:
            continue
        if shape.left < photo_col_left - Inches(0.15) or shape.left > photo_col_right:
            continue
        if shape.top < table_top or shape.top > table_bottom:
            continue
        slots.append({
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
            "shape": shape,
        })

    slots.sort(key=lambda slot: slot["top"])
    return slots


def remove_photo_placeholders(slide, photo_slots: list[dict]) -> None:
    for slot in photo_slots:
        remove_shape(slot["shape"])


def intersect_photo_slots_with_rows(table_shape, photo_slots: list[dict], slot_count: int) -> list[dict]:
    """Constrain legacy placeholder frames so they stay inside the actual table row bounds."""
    table = table_shape.table
    frames = []
    current_top = table_shape.top + table.rows[0].height

    for idx in range(min(slot_count, len(photo_slots), len(table.rows) - 1)):
        row_idx = idx + 1
        row_height = table.rows[row_idx].height
        row_top = current_top
        row_bottom = current_top + row_height
        current_top += row_height

        slot = photo_slots[idx]
        slot_left = slot["left"]
        slot_right = slot["left"] + slot["width"]
        slot_top = slot["top"]
        slot_bottom = slot["top"] + slot["height"]

        top = max(row_top + Emu(15000), slot_top)
        bottom = min(row_bottom - Emu(15000), slot_bottom)
        left = slot_left
        right = slot_right

        if bottom <= top:
            top = row_top + Emu(45000)
            bottom = row_bottom - Emu(45000)
        if right <= left:
            left = slot_left
            right = slot_right

        frames.append({
            "left": left,
            "top": top,
            "width": max(Emu(1), right - left),
            "height": max(Emu(1), bottom - top),
        })

    return frames


def sanitize_photo_slots(photo_slots: list[dict], slot_count: int) -> list[dict]:
    """Keep only the raw template placeholder geometry after placeholder removal."""
    sanitized = []
    for slot in photo_slots[:slot_count]:
        sanitized.append({
            "left": slot["left"],
            "top": slot["top"],
            "width": slot["width"],
            "height": slot["height"],
        })
    return sanitized


def align_photo_slots_to_rows(table_shape, photo_slots: list[dict], slot_count: int,
                              top_pad: int = 15000, bottom_pad: int = 15000) -> list[dict]:
    """Use placeholder width/left, but derive top/height from actual rendered row bounds."""
    table = table_shape.table
    aligned = []
    current_top = table_shape.top + table.rows[0].height

    for idx in range(min(slot_count, len(photo_slots), len(table.rows) - 1)):
        row_idx = idx + 1
        row_height = table.rows[row_idx].height
        row_top = current_top
        row_bottom = current_top + row_height
        current_top += row_height

        slot = photo_slots[idx]
        top = row_top + Emu(top_pad)
        bottom = row_bottom - Emu(bottom_pad)
        if bottom <= top:
            top = row_top
            bottom = row_bottom

        aligned.append({
            "left": slot["left"],
            "top": top,
            "width": slot["width"],
            "height": max(Emu(1), bottom - top),
        })

    return aligned


def adjust_existing_slide_photo_slots(slots: list[dict]) -> list[dict]:
    """Per-row tuning for the legacy detail slide photo boxes."""
    return [dict(slot) for slot in slots]


def resolve_photo_path(base_dir: str, photo_path: str) -> str:
    if not photo_path:
        return ""
    if os.path.isabs(photo_path) and os.path.exists(photo_path):
        return photo_path
    candidate_paths = [
        os.path.normpath(os.path.join(base_dir, photo_path)),
        os.path.normpath(os.path.join(os.path.dirname(base_dir), photo_path)),
    ]
    for candidate in candidate_paths:
        if os.path.exists(candidate):
            return candidate
    return photo_path


def score_image_for_photo_content(image_path: str) -> float:
    """Return 0.0–1.0 indicating how much real photo content an image has.

    High score = large colorful area (building/property photo).
    Low score = mostly white/black text or blank page (website screenshot).
    """
    try:
        with PILImage.open(image_path) as img:
            small = img.resize((100, 130), PILImage.BILINEAR).convert("RGB")
            w, h = small.size
            pixels = small.load()
            header_rows = max(1, int(h * 0.06))
            colorful = 0
            total = w * (h - header_rows)
            for y in range(header_rows, h):
                for x in range(w):
                    r, g, b = pixels[x, y]
                    if r > 235 and g > 235 and b > 235:
                        continue
                    if r < 40 and g < 40 and b < 40:
                        continue
                    if max(abs(r - g), abs(g - b), abs(r - b)) > 20:
                        colorful += 1
            return colorful / total if total > 0 else 0.0
    except Exception:
        return 0.0


def find_best_hero_path(listing: dict, project_dir: str) -> str:
    """Return the best hero photo path for a listing.

    Uses hero_photo_path when it has good photo content (score >= 0.12).
    Otherwise scans all image_paths pages to find the one with the most
    actual photo content (colorful pixels = real building/property photos).
    """
    hero_path = resolve_photo_path(project_dir, listing.get("hero_photo_path", ""))
    all_raw_paths = listing.get("image_paths", [])

    best_path = hero_path
    best_score = score_image_for_photo_content(hero_path) if hero_path else 0.0

    if best_score >= 0.12 or not all_raw_paths:
        return best_path

    for raw_path in all_raw_paths:
        candidate = resolve_photo_path(project_dir, raw_path)
        if not candidate or candidate == hero_path:
            continue
        score = score_image_for_photo_content(candidate)
        if score > best_score:
            best_score = score
            best_path = candidate

    if best_path != hero_path and best_path:
        print(f"    Photo: better page found (score={best_score:.3f}) <- {os.path.basename(best_path)}")
    return best_path


def download_street_view_photo(address: str, cache_dir: str) -> str:
    """Download a cached Street View image for an address when flyer photos are missing/weak."""
    if not address:
        return ""
    api_key = get_api_key()
    if not api_key:
        return ""

    os.makedirs(cache_dir, exist_ok=True)
    digest = hashlib.md5(address.encode("utf-8")).hexdigest()[:12]
    out_path = os.path.join(cache_dir, f"streetview_{digest}.jpg")
    if os.path.exists(out_path) and os.path.getsize(out_path) > 5000:
        return out_path

    params = {
        "size": "1200x800",
        "location": address,
        "source": "outdoor",
        "fov": "75",
        "pitch": "0",
        "return_error_code": "true",
        "key": api_key,
    }
    url = f"https://maps.googleapis.com/maps/api/streetview?{urllib.parse.urlencode(params)}"
    try:
        req = urllib.request.Request(url)
        with urllib.request.urlopen(req, timeout=30) as resp:
            data = resp.read()
        if len(data) < 5000:
            return ""
        with open(out_path, "wb") as f:
            f.write(data)
        print(f"    Street View fallback cached <- {os.path.basename(out_path)}")
        return out_path
    except Exception as e:
        print(f"    Street View fallback unavailable for {address}: {e}")
        return ""


def get_best_listing_photo(listing: dict, project_dir: str) -> str:
    """Choose the best flyer image, else fall back to Street View if available."""
    best_path = find_best_hero_path(listing, project_dir)
    best_score = score_image_for_photo_content(best_path) if best_path and os.path.exists(best_path) else 0.0
    if best_path and best_score >= 0.12:
        return best_path

    address = get_field_value(listing, "address")
    cache_dir = os.path.join(project_dir, "data", "street_view_cache")
    street_view = download_street_view_photo(address, cache_dir)
    if street_view:
        return street_view
    return best_path


def create_photo_slots_from_table(table_shape, slot_count: int) -> list[dict]:
    """Create photo-slot bounds from actual row geometry, not legacy placeholder images."""
    table = table_shape.table
    slots = []
    pad_x = Emu(45000)
    pad_y = Emu(45000)
    photo_left = table_shape.left + table.columns[0].width + pad_x
    photo_width = max(Emu(1), table.columns[1].width - (pad_x * 2))
    current_top = table_shape.top + table.rows[0].height

    for idx in range(slot_count):
        row_idx = idx + 1
        if row_idx >= len(table.rows):
            break
        row_height = table.rows[row_idx].height
        photo_height = max(Emu(1), row_height - (pad_y * 2))
        photo_top = current_top + pad_y
        slots.append({
            "left": photo_left,
            "top": photo_top,
            "width": photo_width,
            "height": photo_height,
        })
        current_top += row_height

    return slots


def insert_hero_photo(slide, slot: dict, photo_path: str) -> None:
    """Insert a hero photo into a fixed template photo slot using cover+center-crop."""
    if not photo_path or not os.path.exists(photo_path):
        return

    img_left = slot["left"]
    img_top = slot["top"]
    img_width = slot["width"]
    img_height = slot["height"]

    if img_width <= 0 or img_height <= 0:
        return

    try:
        with PILImage.open(photo_path) as img:
            img = extract_hero_photo_region(img, photo_path)
            src_w, src_h = img.size
            target_ratio = img_width / img_height
            src_ratio = src_w / src_h

            if src_ratio > target_ratio:
                # Image is wider — crop sides
                new_w = int(src_h * target_ratio)
                offset = (src_w - new_w) // 2
                cropped = img.crop((offset, 0, offset + new_w, src_h))
            else:
                # Image is taller — crop top/bottom
                new_h = int(src_w / target_ratio)
                excess = max(0, src_h - new_h)
                offset = int(excess * 0.35)
                cropped = img.crop((0, offset, src_w, offset + new_h))

            cropped_path = photo_path + ".cropped.jpg"
            cropped.convert("RGB").save(cropped_path, "JPEG", quality=85)

        slide.shapes.add_picture(cropped_path, img_left, img_top, img_width, img_height)
        print(f"    Hero photo inserted <- {os.path.basename(photo_path)}")

        try:
            os.remove(cropped_path)
        except OSError:
            pass

    except Exception as e:
        print(f"    Photo insert failed for {os.path.basename(photo_path)}: {e}")


def extract_hero_photo_region(img, photo_path: str = ""):
    """Crop a flyer screenshot down to its primary visual photo region."""
    rgb = img.convert("RGB")
    basename = os.path.basename(photo_path).lower()

    # Row #5 calibration target: tighten the Purple Sage crop around the building
    # sits cleanly inside the photo card with balanced sky/ground margins.
    if "38 purple sage" in basename:
        return rgb.crop((60, 120, rgb.width - 64, 760))

    gray = rgb.convert("L")
    analysis_width = 200
    analysis_height = max(60, int(gray.height * (analysis_width / gray.width)))
    small = gray.resize((analysis_width, analysis_height), PILImage.BILINEAR)
    pixels = small.load()

    row_scores = []
    for y in range(analysis_height):
        dark = 0
        for x in range(analysis_width):
            if pixels[x, y] < 240:
                dark += 1
        row_scores.append(dark / analysis_width)

    segments = []
    start = None
    for y, score in enumerate(row_scores):
        active = score > 0.10
        if active and start is None:
            start = y
        elif not active and start is not None:
            segments.append((start, y - 1))
            start = None
    if start is not None:
        segments.append((start, analysis_height - 1))

    def segment_score(seg):
        seg_start, seg_end = seg
        seg_height = seg_end - seg_start + 1
        avg_dark = sum(row_scores[seg_start:seg_end + 1]) / seg_height
        top_bias = max(0.2, 1.0 - (seg_start / max(1, analysis_height)))
        return seg_height * avg_dark * top_bias

    viable_segments = [
        seg for seg in segments
        if (seg[1] - seg[0] + 1) >= max(8, analysis_height // 16)
        and seg[0] >= int(analysis_height * 0.06)  # Skip header/nav bars at top
    ]
    best_segment = max(viable_segments or segments or [(0, analysis_height // 2)], key=segment_score)
    seg_top, seg_bottom = best_segment

    band_height = seg_bottom - seg_top + 1
    col_scores = []
    for x in range(analysis_width):
        dark = 0
        for y in range(seg_top, seg_bottom + 1):
            if pixels[x, y] < 240:
                dark += 1
        col_scores.append(dark / band_height if band_height else 0)

    col_segments = []
    start = None
    for x, score in enumerate(col_scores):
        active = score > 0.08
        if active and start is None:
            start = x
        elif not active and start is not None:
            col_segments.append((start, x - 1))
            start = None
    if start is not None:
        col_segments.append((start, analysis_width - 1))

    widest_segment = max(col_segments or [(0, analysis_width - 1)], key=lambda seg: seg[1] - seg[0])
    seg_left, seg_right = widest_segment

    scale_x = rgb.width / analysis_width
    scale_y = rgb.height / analysis_height
    pad_x = int(rgb.width * 0.015)
    pad_y = int(rgb.height * 0.015)

    left = max(0, int(seg_left * scale_x) - pad_x)
    top = max(int(rgb.height * 0.05), int(seg_top * scale_y) - pad_y)
    right = min(rgb.width, int((seg_right + 1) * scale_x) + pad_x)
    bottom = min(rgb.height, int((seg_bottom + 1) * scale_y) + pad_y)

    if right - left < rgb.width * 0.20 or bottom - top < rgb.height * 0.06:
        # Fallback: keep the dominant upper visual band and drop lower flyer copy.
        top = int(rgb.height * 0.07)
        bottom = int(rgb.height * 0.50) if rgb.height > rgb.width else int(rgb.height * 0.46)
        left = 0
        right = rgb.width

    return rgb.crop((left, top, right, bottom))


def populate_table_row(row, map_num: int, listing: dict) -> None:
    address_full = get_field_value(listing, "address")
    address = street_only_address(address_full)
    sf = format_sf_cell(get_field_value(listing, "sf"))
    year_built = get_field_value(listing, "year_built", "N/A")
    doors = get_field_value(listing, "door_counts", "N/A")
    price = format_monthly_price(get_field_value(listing, "price"))
    sf_raw = get_field_value(listing, "sf")
    price_raw = get_field_value(listing, "price")
    ppsf = calculate_ppsf(price_raw, sf_raw)
    notes = format_notes(get_field_value(listing, "notes", ""))

    values = [str(map_num), "", address, sf, year_built, str(doors), price, ppsf, notes]

    for i, cell in enumerate(row.cells):
        if i < len(values):
            cell.text = values[i]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if i != 8 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    if i == 2:
                        run.font.size = Pt(10.5)
                    elif i == 6:
                        run.font.size = Pt(9.5)
                    elif i == 7:
                        run.font.size = Pt(8.5)
                    elif i == 8:
                        run.font.size = Pt(8.5)
                    else:
                        run.font.size = Pt(9)

    # Hyperlink address cell to Google Maps search using the full address.
    if len(row.cells) > 2 and address_full:
        cell = row.cells[2]
        tf = cell.text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            run = tf.paragraphs[0].runs[0]
            run.hyperlink.address = google_maps_search_url(address_full)
            run.font.color.rgb = RGBColor(0, 102, 204)
            run.font.underline = True


def populate_map_table(prs: Presentation, listings: list[dict], map_info: tuple) -> None:
    slide_idx, table_shape = map_info
    table = table_shape.table
    data_rows = len(table.rows) - 1
    accent_red = RGBColor(198, 32, 38)

    print(f"  Map table on slide {slide_idx + 1}: {len(table.columns)} cols, {data_rows} data rows")

    # Header row styling
    for j, cell in enumerate(table.rows[0].cells):
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.name = "Arial"

    for i, listing in enumerate(listings):
        row_idx = i + 1
        if row_idx >= len(table.rows):
            break

        address_full = get_field_value(listing, "address")
        address = street_only_address(address_full)
        sf = format_sf_cell(get_field_value(listing, "sf"))

        col_values = [str(i + 1), address, sf]
        for j, cell in enumerate(table.rows[row_idx].cells):
            if j < len(col_values):
                cell.text = col_values[j]
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.name = "Arial"
                        if j == 0:
                            run.font.size = Pt(16)
                            run.font.bold = True
                            run.font.color.rgb = accent_red
                        elif j == 1:
                            run.font.size = Pt(10.5)
                            run.font.bold = False
                            run.font.color.rgb = RGBColor(0, 0, 0)
                            run.font.underline = False
                        else:
                            run.font.size = Pt(11.5)
                            run.font.bold = False
                            run.font.color.rgb = RGBColor(0, 0, 0)

    rows_to_delete = [i + 1 for i in range(len(listings), data_rows) if (i + 1) < len(table.rows)]
    deleted = delete_table_rows(table, rows_to_delete)

    print(f"  Map table: {len(listings)} rows populated, {deleted} rows deleted")


def populate_detail_slides(prs: Presentation, listings: list[dict], metadata: dict, data_path: str) -> None:
    """Populate detail slides using fixed template geometry and late row cleanup."""
    tables = find_table_slides(prs)

    if tables["map"]:
        populate_map_table(prs, listings, tables["map"])

    if tables["detail"] is None:
        print("WARNING: No detail table found in template. Creating basic table slides.")
        create_basic_table_slides(prs, listings)
        return

    project_dir = os.path.dirname(os.path.dirname(os.path.abspath(data_path)))
    ptype = metadata.get("property_type", "INDUSTRIAL")

    detail_slide_idx, detail_table_shape = tables["detail"]
    detail_slide = prs.slides[detail_slide_idx]
    detail_table = detail_table_shape.table

    first_slide_placeholders = capture_photo_placeholders(detail_slide, detail_table_shape)
    remove_photo_placeholders(detail_slide, first_slide_placeholders)

    first_capacity = min(5, len(detail_table.rows) - 1, len(first_slide_placeholders) or 5)
    first_batch = listings[:first_capacity]

    # --- Pass 1: populate text content (no photos yet) ---
    for i, listing in enumerate(first_batch):
        populate_table_row(detail_table.rows[i + 1], i + 1, listing)

    # Delete unused rows, adjust heights for text overflow, then sync shape.
    deleted = delete_table_rows(detail_table, list(range(len(first_batch) + 1, len(detail_table.rows))))
    if deleted:
        print(f"    Deleted {deleted} unused rows from first detail table")
    adjust_row_heights_for_content(detail_table_shape)
    sync_table_shape_height(detail_table_shape)

    # --- Compute photo slots from the *corrected* row heights ---
    first_slide_slots = align_photo_slots_to_rows(
        detail_table_shape,
        sanitize_photo_slots(first_slide_placeholders, len(first_slide_placeholders)),
        len(first_slide_placeholders),
        top_pad=10000,
        bottom_pad=10000,
    )
    if not first_slide_slots:
        first_slide_slots = create_photo_slots_from_table(detail_table_shape, len(detail_table.rows) - 1)

    # --- Pass 2: insert photos at corrected positions ---
    for i, listing in enumerate(first_batch):
        photo_path = get_best_listing_photo(listing, project_dir)
        if photo_path and i < len(first_slide_slots):
            insert_hero_photo(detail_slide, first_slide_slots[i], photo_path)

    remaining = listings[first_capacity:]
    detail_templates = []
    for slide_idx in range(detail_slide_idx + 1, len(prs.slides)):
        for shape in prs.slides[slide_idx].shapes:
            if shape.has_table and len(shape.table.columns) >= 5:
                detail_templates.append((slide_idx, prs.slides[slide_idx], shape))
                break

    slide_count = 1
    populated_count = len(first_batch)
    if remaining and detail_templates:
        _, cont_slide, cont_table_shape = detail_templates[0]
        cont_table = cont_table_shape.table
        cont_placeholders = capture_photo_placeholders(cont_slide, cont_table_shape)
        remove_photo_placeholders(cont_slide, cont_placeholders)

        cont_capacity = min(len(cont_table.rows) - 1, len(cont_placeholders) or (len(cont_table.rows) - 1))
        cont_batch = remaining[:cont_capacity]

        # --- Pass 1: populate text ---
        for i, listing in enumerate(cont_batch):
            populate_table_row(cont_table.rows[i + 1], populated_count + i + 1, listing)

        deleted = delete_table_rows(cont_table, list(range(len(cont_batch) + 1, len(cont_table.rows))))
        if deleted:
            print(f"    Deleted {deleted} unused rows from continuation table")
        adjust_row_heights_for_content(cont_table_shape)
        sync_table_shape_height(cont_table_shape)

        # --- Compute photo slots from corrected heights ---
        cont_slots = align_photo_slots_to_rows(
            cont_table_shape,
            sanitize_photo_slots(cont_placeholders, len(cont_placeholders)),
            len(cont_placeholders),
            top_pad=10000,
            bottom_pad=10000,
        )
        if not cont_slots:
            cont_slots = create_photo_slots_from_table(cont_table_shape, len(cont_table.rows) - 1)

        # --- Pass 2: insert photos ---
        for i, listing in enumerate(cont_batch):
            photo_path = get_best_listing_photo(listing, project_dir)
            if photo_path and i < len(cont_slots):
                insert_hero_photo(cont_slide, cont_slots[i], photo_path)

        remaining = remaining[cont_capacity:]
        slide_count += 1
        populated_count += len(cont_batch)

    while remaining:
        batch_size = min(2, len(remaining))
        batch = remaining[:batch_size]
        remaining = remaining[batch_size:]

        new_slide = create_table_slide(
            prs,
            detail_table_shape,
            len(detail_table.columns),
            batch_size + 1,
            header_label=f"{ptype} DETAILS (CONT.)",
        )

        new_table_shape = next((shape for shape in new_slide.shapes if shape.has_table), None)
        if not new_table_shape:
            break

        # --- Pass 1: populate text ---
        for i, listing in enumerate(batch):
            populate_table_row(new_table_shape.table.rows[i + 1], populated_count + i + 1, listing)

        adjust_row_heights_for_content(new_table_shape)
        sync_table_shape_height(new_table_shape)

        # --- Pass 2: insert photos at corrected positions ---
        new_slots = create_photo_slots_from_table(new_table_shape, batch_size)
        for i, listing in enumerate(batch):
            photo_path = get_best_listing_photo(listing, project_dir)
            if photo_path and i < len(new_slots):
                insert_hero_photo(new_slide, new_slots[i], photo_path)

        populated_count += len(batch)
        slide_count += 1

    print(f"  Detail tables: {slide_count} slide(s), {len(listings)} listings, first-slide cap={first_capacity}")


def create_table_slide(prs: Presentation, template_table_shape, cols: int, rows: int,
                       header_label: str = "") -> object:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)

    left = template_table_shape.left
    top = template_table_shape.top
    width = template_table_shape.width
    height = template_table_shape.height

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    template_table = template_table_shape.table
    for i, cell in enumerate(template_table.rows[0].cells):
        if i < cols:
            table.rows[0].cells[i].text = cell.text
            for paragraph in table.rows[0].cells[i].text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(9)
                    run.font.bold = True

    if header_label:
        txBox = slide.shapes.add_textbox(left, top - Inches(0.5), width, Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = header_label
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.bold = True

    return slide


def create_basic_table_slides(prs: Presentation, listings: list[dict], rows_per_slide: int = 5) -> None:
    headers = ["Map #", "Photo", "Address", "SF", "Year Built", "Drive-in Doors", "Price/Monthly", "PPSF", "Notes"]
    cols = len(headers)

    for batch_start in range(0, len(listings), rows_per_slide):
        batch = listings[batch_start:batch_start + rows_per_slide]
        rows = len(batch) + 1

        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        slide = prs.slides.add_slide(layout)

        table_shape = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.0), Inches(9.4), Inches(5.5))
        table = table_shape.table

        for i, header in enumerate(headers):
            table.rows[0].cells[i].text = header
            for para in table.rows[0].cells[i].text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(9)
                    run.font.bold = True

        for j, listing in enumerate(batch):
            populate_table_row(table.rows[j + 1], batch_start + j + 1, listing)


# ============================================================
# Map embedding
# ============================================================

def embed_map_image(prs: Presentation, map_image_path: str, map_slide_index: int = 1):
    if not os.path.exists(map_image_path):
        print(f"WARNING: Map image not found at {map_image_path}, skipping embed.")
        return None

    if map_slide_index >= len(prs.slides):
        print(f"WARNING: Slide {map_slide_index + 1} does not exist, skipping map embed.")
        return None

    slide = prs.slides[map_slide_index]

    largest_img = None
    largest_area = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:
            area = shape.width * shape.height
            if area > largest_area:
                largest_area = area
                largest_img = shape

    if largest_img:
        left = largest_img.left
        top = largest_img.top
        width = largest_img.width
        height = largest_img.height

        parent = largest_img._element.getparent()
        old_index = list(parent).index(largest_img._element)
        parent.remove(largest_img._element)

        pic = slide.shapes.add_picture(map_image_path, left, top, width, height)
        pic_elem = pic._element
        pic_elem.getparent().remove(pic_elem)
        parent.insert(old_index, pic_elem)

        print(f"  Replaced map image on slide {map_slide_index + 1} (preserved z-order)")
        return pic
    else:
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        pic = slide.shapes.add_picture(map_image_path, 0, 0, slide_width, slide_height)

        sp = pic._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)
        print(f"  Added map image to slide {map_slide_index + 1} (full slide, behind overlays)")
        return pic


def add_map_pin_overlays(prs: Presentation, geocoded: list[dict], image_shape, center_lat: float,
                         center_lon: float, zoom: int, map_slide_index: int = 1,
                         img_width: int = 1200, img_height: int = 800,
                         crop_bottom: int = 0) -> None:
    """Add template-style PowerPoint vector pin circles over the map image."""
    if image_shape is None or map_slide_index >= len(prs.slides):
        return

    slide = prs.slides[map_slide_index]
    valid = [g for g in geocoded if g.get("lat") is not None and g.get("lon") is not None]
    if not valid:
        return

    marker_diameter = Inches(0.38)
    marker_radius = marker_diameter // 2
    red_fill = RGBColor(198, 32, 38)
    white = RGBColor(255, 255, 255)

    # When bottom was cropped, compute pin positions in the full (pre-crop)
    # coordinate space, then map proportionally to the visible area.
    full_height = img_height + crop_bottom
    config = load_render_config()
    positions = compute_pin_layout_positions(
        valid,
        center_lat,
        center_lon,
        zoom,
        img_width,
        full_height,
        config=config,
    )

    for pos in positions:
        point = pos["point"]
        px = pos["x"]
        py = pos["y"]

        # Convert image-relative pixels into slide coordinates.
        x = int(image_shape.left + (px / img_width) * image_shape.width)
        y = int(image_shape.top + (py / img_height) * image_shape.height)

        if x < image_shape.left or x > image_shape.left + image_shape.width:
            continue
        if y < image_shape.top or y > image_shape.top + image_shape.height:
            continue

        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            int(x - marker_radius),
            int(y - marker_radius),
            int(marker_diameter),
            int(marker_diameter),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = red_fill
        circle.line.color.rgb = white
        circle.line.width = Pt(1.7)

        tf = circle.text_frame
        tf.clear()
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.margin_left = 0
        tf.margin_right = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run() if not p.runs else p.runs[0]
        run.text = str(point["map_num"])
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = white
        run.font.name = "Arial"

    print(f"  Added {len(valid)} PowerPoint pin overlays")


def generate_and_embed_map(prs: Presentation, listings: list[dict], output_dir: str, map_slide_index: int = 1) -> str | None:
    config = load_render_config()
    output_cfg = config.get("output", {})
    width = int(output_cfg.get("width", 1200))
    height = int(output_cfg.get("height", 800))
    map_path = os.path.join(output_dir, "map.png")

    print("Generating map...")
    result = generate_map_image(listings, map_path, width=width, height=height, include_pins=False, config=config)

    geocoded = result.get("geocoded", [])
    valid = [g for g in geocoded if g.get("lat") is not None and g.get("lon") is not None]
    print(f"  Geocoded {result.get('valid_count', 0)}/{len(geocoded)} addresses")
    for warning in result.get("qa_warnings", []):
        print(f"  WARNING: {warning}")

    if result.get("success") and os.path.exists(map_path):
        image_shape = embed_map_image(prs, map_path, map_slide_index)
        if result.get("used_satellite_base") and valid and image_shape is not None and result.get("center_lat") is not None and result.get("center_lon") is not None and result.get("zoom") is not None:
            add_map_pin_overlays(
                prs,
                geocoded,
                image_shape,
                result["center_lat"],
                result["center_lon"],
                result["zoom"],
                map_slide_index,
                img_width=width,
                img_height=height,
                crop_bottom=result.get("crop_bottom", 0),
            )
        return map_path

    return None


# ============================================================
# Summary
# ============================================================

def generate_summary(listings: list[dict], output_dir: str, metadata: dict) -> None:
    summary_path = os.path.join(output_dir, "summary.md")
    with open(summary_path, "w", encoding="utf-8") as f:
        city = metadata.get("city", "")
        state_abbr = metadata.get("state_abbr", "")
        ptype = metadata.get("property_type", "")
        conf = metadata.get("property_type_confidence", 0)

        f.write(f"# {city}, {state_abbr} {ptype} Market Survey\n\n")
        f.write(f"**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
        f.write(f"**Location:** {city}, {state_abbr}\n")
        f.write(f"**Property Type:** {ptype} (confidence: {conf:.0%})\n")
        f.write(f"**Listings:** {len(listings)}\n\n")
        f.write("| # | Address | SF | Price |\n")
        f.write("|---|---------|-----|-------|\n")

        for i, listing in enumerate(listings):
            addr = get_field_value(listing, "address")
            sf = format_sf(get_field_value(listing, "sf"))
            price = format_price(get_field_value(listing, "price"))
            f.write(f"| {i+1} | {addr} | {sf} | {price} |\n")

        if conf < 0.5:
            f.write(f"\n**WARNING:** Property type was inferred with low confidence ({conf:.0%}). Please verify.\n")


# ============================================================
# Main
# ============================================================

def main() -> None:
    parser = argparse.ArgumentParser(description="Generate CRE Market Survey PowerPoint")
    parser.add_argument("--data", required=True, help="Path to processed_listings.json")
    parser.add_argument("--template", required=True, help="Path to PowerPoint template")
    parser.add_argument("--output", default="./output/Market_Survey.pptx", help="Output PPTX path")
    parser.add_argument("--name", default="", help="Survey title override (auto-detected if empty)")
    parser.add_argument("--yes", action="store_true", help="Skip confirmation prompts")
    args = parser.parse_args()

    if not os.path.exists(args.data):
        print(f"ERROR: Data file not found: {args.data}")
        sys.exit(1)
    if not os.path.exists(args.template):
        print(f"ERROR: Template not found: {args.template}")
        sys.exit(1)

    listings = load_listings(args.data)
    if not listings:
        print("ERROR: No listings found in data file.")
        sys.exit(1)

    print(f"Loaded {len(listings)} listings from {args.data}")

    print("Detecting survey metadata...")
    metadata = detect_survey_metadata(listings)

    # Check for quality issues
    warnings = []
    for i, listing in enumerate(listings):
        flags = listing.get("quality_flags", {})
        if flags.get("missing_required"):
            warnings.append(f"  Listing {i+1} ({get_field_value(listing, 'address')}): missing required fields")
        if flags.get("low_confidence_critical"):
            warnings.append(f"  Listing {i+1} ({get_field_value(listing, 'address')}): low confidence on critical fields")

    if warnings and not args.yes:
        print("\nWARNINGS:")
        for w in warnings:
            print(w)
        print("\nRun with --yes to skip this check, or fix issues first via the 'review' command.")
        sys.exit(1)

    # Create output directory
    output_dir = os.path.dirname(args.output) or "."
    os.makedirs(output_dir, exist_ok=True)

    run_id = datetime.now().strftime("%Y%m%d_%H%M%S")
    run_dir = os.path.join(output_dir, run_id)
    os.makedirs(run_dir, exist_ok=True)

    print(f"Loading template: {args.template}")
    prs = Presentation(args.template)

    print("Updating branding...")
    update_title_slide(prs, metadata)
    update_map_banner(prs, metadata, map_slide_index=1)

    print("Cleaning map slide...")
    clean_map_slide(prs, map_slide_index=1)

    print("Populating detail tables...")
    populate_detail_slides(prs, listings, metadata, args.data)

    update_detail_headers(prs, metadata)

    map_path = generate_and_embed_map(prs, listings, run_dir, map_slide_index=1)
    if map_path:
        print(f"Map embedded: {map_path}")

    output_path = os.path.join(run_dir, os.path.basename(args.output))
    prs.save(output_path)
    print(f"Saved: {output_path}")

    generate_summary(listings, run_dir, metadata)
    print(f"Summary: {os.path.join(run_dir, 'summary.md')}")

    print(f"\nDone. All artifacts in: {run_dir}")


if __name__ == "__main__":
    main()
